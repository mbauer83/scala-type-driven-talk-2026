#!/usr/bin/env python3
"""Build talk.pptx with the speaker notes as REAL TEXT in the notes pane.

Why this exists
---------------
pympress and pdfpc live in WSL2 and are not usable for presenting on Windows.
`touying compile --format pptx` produces slide images and **no notes at all**,
so PowerPoint's presenter view has nothing to show — which is the whole point of
using it. This script renders the same pages and puts each slide's script into
the PPTX notes field, so on Windows you get presenter view, notes, timer and the
next-slide thumbnail with no WSL involved.

Why the slides are still images
-------------------------------
There is no lossless Typst-to-OOXML path. Typst has no vector export to the
PowerPoint shape model, and PDF-to-PPTX converters reconstruct text as hundreds
of positioned boxes — which loses the two custom fonts, the code-pane tints and
the absolute layout, and would need all 49 slides re-checked by eye. The slides
are not meant to be edited in PowerPoint; the Typst source is the truth. What
has to survive the trip is the NOTES, and those are real text here.

Rendered at 192 ppi (2560x1440) so a 4K projector still gets clean type.

    python3 tools/make-pptx.py [--ppi 192] [--full-notes] [--output talk.pptx]

--full-notes keeps the PREPARATION block. The default drops it: every script
marks that block "not for the night", and PowerPoint's notes pane is small
enough that it would bury the script under its own commentary.
"""
import argparse
import copy
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile

ROOT = os.path.normpath(os.path.join(os.path.dirname(os.path.abspath(__file__)), ".."))
DECK = os.path.join("touying", "deck.typ")
# Vendored fonts — see the note in touying/theme.typ. Keep in step with the
# TYPST_FONTS variable in the Makefile.
FONT_ARGS = ["--font-path", os.path.join("touying", "fonts"), "--ignore-system-fonts"]
PDFPC = os.path.join(ROOT, "talk.pdfpc")

# Scripts are TALKING POINTS -> VERBATIM -> PREPARATION, separated by a rule of
# '=' characters. PREPARATION is explicitly not for the night.
PREP_RULE = re.compile(r"^={20,}\s*$", re.M)


def notes_for_pages():
    """Speaker notes, one per page, from the pdfpc sidecar — one source of truth."""
    if not os.path.exists(PDFPC):
        sys.exit("talk.pdfpc not found — run `make all` first.")
    with open(PDFPC, encoding="utf-8") as fh:
        return [p.get("note") or "" for p in json.load(fh)["pages"]]


def trim(note, keep_prep):
    if keep_prep:
        return note.strip()
    return PREP_RULE.split(note, maxsplit=1)[0].strip()


def render(ppi, outdir):
    """One PNG per page. --root . because slides #read() files above touying/."""
    pattern = os.path.join(outdir, "slide-{0p}.png")
    subprocess.run(
        ["typst", "compile", "--root", ".", *FONT_ARGS, DECK, pattern,
         "--format", "png", "--ppi", str(ppi)],
        cwd=ROOT, check=True,
    )
    return sorted(
        os.path.join(outdir, f) for f in os.listdir(outdir) if f.endswith(".png")
    )


def write_note(text_frame, note):
    """Lay a note out so it can be read at a glance in a presenter pane.

    `TextFrame.text = ...` starts a new PARAGRAPH at every line feed, and PPTX
    paragraphs carry no space before or after, so a blank line in the script came
    out as nothing at all — every line the same distance apart, which is what made
    the notes unreadable in Impress.

    So: split the note into blocks on blank lines, one paragraph per block, with
    real space after it. Lines inside a block are line breaks, which keeps the
    runbook's aligned columns and the FILE / DIR / COMMAND header intact.
    """
    from pptx.util import Pt
    from pptx.oxml.ns import qn

    blocks = [b for b in re.split(r"\n\s*\n", note.strip()) if b.strip()]
    text_frame.word_wrap = True
    first = True
    for block in blocks:
        # A SPACER PARAGRAPH, not just space-after. LibreOffice Impress ignores
        # spcAft on notes text and collapses a genuinely empty paragraph, which
        # is why the blank lines in the script disappeared entirely. A paragraph
        # holding one space cannot be collapsed by anything, so the gap survives
        # every renderer. spcAft stays as well — PowerPoint does honour it.
        if not first:
            gap = text_frame.add_paragraph()
            grun = gap.add_run()
            grun.text = " "
            grun.font.size = Pt(8)
        para = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        lines = block.split("\n")
        # A block whose lines are aligned into columns — the FILE / DIR / COMMAND
        # header and the runbook table — only reads correctly in a fixed pitch.
        # Courier New rather than the deck's JetBrains Mono: PPTX cannot carry a
        # vendored font, and Courier New is present on every machine that opens
        # this file. Prose keeps the reader's default face.
        # Column alignment, not merely indentation: a run of three or more
        # spaces BETWEEN two non-space characters is what padding a column looks
        # like. A hanging indent under a numbered talking point is not that, and
        # reads better in the reader's own face.
        fixed = sum(1 for l in lines if re.search(r"\S {3,}\S", l)) >= 1
        for m, line in enumerate(lines):
            run = para.add_run()
            run.text = line
            run.font.size = Pt(11 if fixed else 12)
            if fixed:
                run.font.name = "Courier New"
            if m < len(lines) - 1:
                # a line break inside the paragraph, not a new paragraph
                br = copy.deepcopy(run._r)
                br.tag = qn("a:br")
                for child in list(br):
                    br.remove(child)
                run._r.addnext(br)
        para.space_after = Pt(10)
        para.line_spacing = 1.15


def build(images, notes, output):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)          # 16:9, the deck's aspect
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]              # no placeholders to fight with

    for i, image in enumerate(images):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(image, 0, 0, prs.slide_width, prs.slide_height)
        note = notes[i] if i < len(notes) else ""
        if note:
            write_note(slide.notes_slide.notes_text_frame, note)

    prs.save(output)
    return len(images)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--ppi", type=int, default=192)
    ap.add_argument("--full-notes", action="store_true")
    ap.add_argument("--output", default=os.path.join(ROOT, "talk.pptx"))
    args = ap.parse_args()

    notes = [trim(n, args.full_notes) for n in notes_for_pages()]
    tmp = tempfile.mkdtemp(prefix="talk-pptx-")
    try:
        images = render(args.ppi, tmp)
        if len(images) != len(notes):
            print(f"warning: {len(images)} pages rendered, {len(notes)} notes in talk.pdfpc",
                  file=sys.stderr)
        n = build(images, notes, args.output)
    finally:
        shutil.rmtree(tmp, ignore_errors=True)

    with_notes = sum(1 for x in notes if x)
    size = os.path.getsize(args.output) / 1e6
    print(f"{args.output}: {n} slides, {with_notes} with notes, {size:.0f} MB, {args.ppi} ppi")


if __name__ == "__main__":
    main()
